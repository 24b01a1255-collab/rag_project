from pptx import Presentation

from langchain_core.documents import Document


def load_pptx(uploaded_file):

    prs = Presentation(uploaded_file)

    all_text = []

    for slide_num, slide in enumerate(prs.slides):

        slide_text = ""

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                slide_text += shape.text + "\n"

        all_text.append(

            Document(

                page_content=slide_text,

                metadata={

                    "source": uploaded_file.name,

                    "page": slide_num + 1
                }
            )
        )

    return all_text